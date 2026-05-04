"""
MixBot – Professional PowerPoint Generator
7 slides, real screenshots, dark glazed design matching app aesthetics.
"""

import os, random
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Palette ───────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x08, 0x10, 0x28)
PANEL       = RGBColor(0x0D, 0x1A, 0x3A)
PANEL2      = RGBColor(0x11, 0x20, 0x45)
GOLD        = RGBColor(0xD4, 0xAF, 0x37)
GOLD_LT     = RGBColor(0xF0, 0xCC, 0x60)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
MUTED       = RGBColor(0xAA, 0xAA, 0xCC)
DIM         = RGBColor(0x66, 0x66, 0x88)
RED         = RGBColor(0xFF, 0x6B, 0x6B)
GREEN       = RGBColor(0x4C, 0xAF, 0x50)
BLUE_ACC    = RGBColor(0x3A, 0x78, 0xCC)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ─── Screenshot paths ───────────────────────────────────────────────────────────
SS_DIR = r"C:\Users\manik\.gemini\antigravity\brain\da806676-745c-4527-95a3-49bf27ac03be"
SS = {
    "landing":    os.path.join(SS_DIR, "ss_01_landing_1777229468860.png"),
    "login":      os.path.join(SS_DIR, "ss_02_login_modal_1777229484942.png"),
    "chat_ui":    os.path.join(SS_DIR, "ss_03_chat_interface_1777229507308.png"),
    "chat_resp":  os.path.join(SS_DIR, "ss_04_chat_response_1777229529607.png"),
    "admin":      os.path.join(SS_DIR, "ss_05_admin_panel_1777229543808.png"),
    "admin_chat": os.path.join(SS_DIR, "ss_06_admin_chat_history_1777229551518.png"),
}

OUTPUT = r"C:\Users\manik\OneDrive\Desktop\Ai Essentials\cocktail\MixBot_Presentation.pptx"

# ─── Primitives ────────────────────────────────────────────────────────────────
def bg(slide, color=NAVY):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def rect(slide, x, y, w, h, fill=PANEL, border=None, bw=Pt(1)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = bw
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h,
        size=14, bold=False, italic=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def hline(slide, x, y, w, color=GOLD, h=Inches(0.035)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()

def vline(slide, x, y, h_val, color=GOLD, w=Inches(0.04)):
    s = slide.shapes.add_shape(1, x, y, w, h_val)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()

def img(slide, key, x, y, w, h):
    path = SS.get(key, "")
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, w, h)
        b = slide.shapes.add_shape(1, x, y, w, h)
        b.fill.background(); b.line.color.rgb = GOLD; b.line.width = Pt(1.2)
    else:
        rect(slide, x, y, w, h, PANEL2, GOLD)
        txt(slide, f"[ {key} ]", x, y + h//2 - Inches(0.2), w, Inches(0.4),
            size=12, color=MUTED, align=PP_ALIGN.CENTER)

def particles(slide, n=14, seed=99):
    random.seed(seed)
    for _ in range(n):
        rx, ry = random.uniform(0.2, 12.8), random.uniform(0.1, 7.1)
        sz = random.uniform(0.04, 0.14)
        d = slide.shapes.add_shape(9, Inches(rx), Inches(ry), Inches(sz), Inches(sz))
        d.fill.solid()
        d.fill.fore_color.rgb = RGBColor(0xD4, 0xAF, 0x37)
        d.line.fill.background()

def header_band(slide, title, subtitle=None):
    rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.5), RGBColor(0x08, 0x14, 0x30))
    hline(slide, Inches(0), Inches(1.5), SLIDE_W)
    txt(slide, title, Inches(0.55), Inches(0.18), Inches(9), Inches(0.8),
        size=32, bold=True, color=GOLD)
    if subtitle:
        txt(slide, subtitle, Inches(0.55), Inches(0.92), Inches(9), Inches(0.5),
            size=13, color=MUTED)

def badge(slide, label, x, y, bg_c=RGBColor(0x18, 0x2E, 0x55), border_c=GOLD):
    bw = Inches(1.85); bh = Inches(0.37)
    r = slide.shapes.add_shape(1, x, y, bw, bh)
    r.fill.solid(); r.fill.fore_color.rgb = bg_c
    r.line.color.rgb = border_c; r.line.width = Pt(0.8)
    tb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.05), bw - Inches(0.2), bh - Inches(0.1))
    tf = tb.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r2 = p.add_run(); r2.text = label
    r2.font.size = Pt(10.5); r2.font.color.rgb = GOLD_LT; r2.font.bold = True

def dot_icon(slide, x, y, color=GOLD):
    d = slide.shapes.add_shape(9, x, y, Inches(0.12), Inches(0.12))
    d.fill.solid(); d.fill.fore_color.rgb = color; d.line.fill.background()

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 – Cover
# ══════════════════════════════════════════════════════════════════════════════
def slide_cover(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl); particles(sl, 22, seed=7)

    # Full-width translucent centre card
    rect(sl, Inches(1.1), Inches(0.55), Inches(11.1), Inches(6.4),
         RGBColor(0x0B, 0x17, 0x35), GOLD, Pt(0.8))

    # Top gold stripe
    hline(sl, Inches(1.1), Inches(0.55), Inches(11.1), GOLD, Inches(0.08))

    # Cocktail hero icon
    txt(sl, "🍹", Inches(5.8), Inches(1.05), Inches(1.8), Inches(1.2),
        size=56, align=PP_ALIGN.CENTER)

    # App name
    txt(sl, "MixBot", Inches(1.5), Inches(2.1), Inches(10.3), Inches(1.3),
        size=78, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Tagline
    txt(sl, "AI-Powered Cocktail Mixologist",
        Inches(2.0), Inches(3.35), Inches(9.3), Inches(0.65),
        size=24, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

    hline(sl, Inches(3.5), Inches(4.05), Inches(6.3), GOLD)

    # Tech pill row
    txt(sl, "Google Gemini 2.0 Flash   ·   FastAPI   ·   OpenRouter   ·   Google OAuth",
        Inches(2.0), Inches(4.2), Inches(9.3), Inches(0.5),
        size=13.5, color=MUTED, align=PP_ALIGN.CENTER)

    # Divider
    hline(sl, Inches(1.1), Inches(6.5), Inches(11.1), RGBColor(0x22, 0x33, 0x55))

    # Footer
    txt(sl, "AI Essentials Capstone Project  |  April 2026",
        Inches(1.5), Inches(6.6), Inches(10.3), Inches(0.4),
        size=11, color=DIM, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 – Project Overview  (Landing + Login screenshot)
# ══════════════════════════════════════════════════════════════════════════════
def slide_overview(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl); particles(sl, 10, seed=12)
    header_band(sl, "Project Overview",
                "What is MixBot and why was it built?")

    # Left: landing screenshot
    img(sl, "landing", Inches(0.4), Inches(1.65), Inches(5.55), Inches(3.9))

    # Right of first screenshot: login screenshot
    img(sl, "login",   Inches(6.15), Inches(1.65), Inches(3.1), Inches(3.9))

    # Right panel: description
    rect(sl, Inches(9.45), Inches(1.65), Inches(3.45), Inches(3.9), PANEL2)
    vline(sl, Inches(9.45), Inches(1.65), Inches(3.9))

    txt(sl, "About the Project", Inches(9.65), Inches(1.82), Inches(3.2), Inches(0.5),
        size=15, bold=True, color=GOLD_LT)
    hline(sl, Inches(9.65), Inches(2.35), Inches(2.9), GOLD, Inches(0.025))

    about = (
        "MixBot is a full-stack AI web app "
        "that acts as your personal bartender. "
        "It uses Gemini 2.0 Flash to understand "
        "natural language and craft unique cocktail "
        "recipes, substitutions and party plans "
        "in real time — all personalised per user."
    )
    txt(sl, about, Inches(9.65), Inches(2.45), Inches(3.15), Inches(2.2),
        size=11.5, color=MUTED, wrap=True)

    highlights = [
        ("🌐", "Web-based chat interface"),
        ("🔐", "JWT + Google OAuth"),
        ("🗄", "SQLite persistent history"),
        ("👑", "Full admin control panel"),
    ]
    y = Inches(4.55)
    for icon, label in highlights:
        dot_icon(sl, Inches(9.65), y + Inches(0.06), GOLD)
        txt(sl, icon + "  " + label, Inches(9.85), y, Inches(2.95), Inches(0.35),
            size=11, color=WHITE)
        y += Inches(0.39)

    # Bottom caption strip
    rect(sl, Inches(0.4), Inches(5.72), Inches(8.3), Inches(0.38), RGBColor(0x0A, 0x14, 0x2E))
    txt(sl, "← Landing hero page (unauthenticated)      ← Login modal with Google OAuth support",
        Inches(0.55), Inches(5.76), Inches(8.1), Inches(0.3),
        size=9.5, color=DIM, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 – Core Features (6 feature cards)
# ══════════════════════════════════════════════════════════════════════════════
def slide_features(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl); particles(sl, 10, seed=33)
    header_band(sl, "Core Features",
                "Everything MixBot can do for you")

    features = [
        ("🍸", "AI Chat Bartender",
         "Powered by Gemini 2.0 Flash. Ask for recipes, substitutions, pairing advice, or party plans in plain English."),
        ("🔐", "Secure Authentication",
         "Email+password login with bcrypt hashing, or instant Google OAuth single sign-on. JWT tokens secure every request."),
        ("📜", "Persistent Chat History",
         "Conversations are saved per-user in SQLite. Sessions sync to the backend on every message — never lose a recipe."),
        ("👑", "Admin Control Panel",
         "Administrators view all users, promote/demote roles, read full chat transcripts, and delete records."),
        ("🛡️", "Rate-Limit Resilience",
         "Exponential back-off retries handle API quota limits. Friendly in-app messages prevent user confusion during outages."),
        ("✨", "Premium Glassmorphism UI",
         "Dark navy canvas, animated gold particles, glass panels, and Playfair Display typography for a luxury feel."),
    ]

    cols, rows = 3, 2
    cw = Inches(4.05); ch = Inches(2.25)
    sx = Inches(0.4);  sy = Inches(1.65)
    gx = Inches(0.22); gy = Inches(0.2)

    for i, (icon, title, desc) in enumerate(features):
        c, r = i % cols, i // cols
        x = sx + c * (cw + gx)
        y = sy + r * (ch + gy)

        rect(sl, x, y, cw, ch, RGBColor(0x0E, 0x1C, 0x40), RGBColor(0x22, 0x38, 0x60), Pt(0.6))
        # Gold left accent strip
        vline(sl, x, y, ch, GOLD)

        txt(sl, icon + "  " + title, x + Inches(0.18), y + Inches(0.16),
            cw - Inches(0.22), Inches(0.45), size=14.5, bold=True, color=GOLD_LT)
        hline(sl, x + Inches(0.18), y + Inches(0.65), cw - Inches(0.35),
              RGBColor(0x2A, 0x3F, 0x70), Inches(0.02))
        txt(sl, desc, x + Inches(0.18), y + Inches(0.75), cw - Inches(0.3), Inches(1.35),
            size=11.5, color=MUTED, wrap=True)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 – Tech Architecture
# ══════════════════════════════════════════════════════════════════════════════
def slide_architecture(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl); particles(sl, 10, seed=55)
    header_band(sl, "System Architecture",
                "How the pieces connect — Frontend → Backend → AI")

    stacks = [
        ("🌐  Frontend", [
            "HTML5 + CSS3 (Vanilla)",
            "Vanilla JavaScript (ES6+)",
            "Glassmorphism + Particle FX",
            "Google Sign-In SDK",
            "Responsive 2-column layout",
        ], RGBColor(0x10, 0x22, 0x48)),
        ("⚙️  Backend", [
            "FastAPI (Python 3.11)",
            "SQLAlchemy ORM",
            "SQLite database",
            "bcrypt password hashing",
            "JWT access tokens",
        ], RGBColor(0x0D, 0x1E, 0x44)),
        ("🤖  AI Engine", [
            "Google Gemini 2.0 Flash",
            "OpenRouter gateway",
            "Custom system prompts",
            "Exponential retry logic",
            "JSON + Chat endpoints",
        ], RGBColor(0x0B, 0x1A, 0x3C)),
    ]

    cw = Inches(3.7);  ch = Inches(4.55)
    sx = Inches(0.55); sy = Inches(1.65); gap = Inches(0.46)

    for i, (title, items, color) in enumerate(stacks):
        x = sx + i * (cw + gap)
        rect(sl, x, sy, cw, ch, color, GOLD, Pt(0.7))
        # thick top bar
        hline(sl, x, sy, cw, GOLD, Inches(0.1))

        txt(sl, title, x + Inches(0.2), sy + Inches(0.18), cw - Inches(0.4), Inches(0.5),
            size=17, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        hline(sl, x + Inches(0.3), sy + Inches(0.75), cw - Inches(0.6),
              RGBColor(0x2A, 0x3E, 0x6A))

        for j, item in enumerate(items):
            dot_icon(sl, x + Inches(0.22), sy + Inches(0.95) + Inches(j * 0.62) + Inches(0.06), GOLD)
            txt(sl, item, x + Inches(0.45), sy + Inches(0.95) + Inches(j * 0.62),
                cw - Inches(0.55), Inches(0.55), size=12.5, color=WHITE)

        # bottom number
        txt(sl, str(i + 1), x + cw - Inches(0.45), sy + ch - Inches(0.42), Inches(0.35), Inches(0.35),
            size=18, bold=True, color=RGBColor(0x2A, 0x3C, 0x66), align=PP_ALIGN.CENTER)

    # Arrow connectors
    for arx in [Inches(4.35), Inches(8.71)]:
        txt(sl, "→", arx, Inches(3.6), Inches(0.4), Inches(0.5),
            size=30, color=GOLD, align=PP_ALIGN.CENTER)

    # Tech badge row at bottom
    badges = ["Python 3.11", "FastAPI", "SQLite", "OpenRouter", "Gemini 2.0 Flash",
              "bcrypt", "JWT", "HTML5 / CSS3"]
    bx = Inches(0.4); by = Inches(6.38)
    for b in badges:
        badge(sl, b, bx, by)
        bx += Inches(1.95)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 – Live Chat Interface (two screenshots side by side)
# ══════════════════════════════════════════════════════════════════════════════
def slide_chat(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl); particles(sl, 10, seed=77)
    header_band(sl, "Live Chat Interface",
                "MixBot responding in real time — admin logged in as 'Admin'")

    # Left: empty interface
    img(sl, "chat_ui",   Inches(0.4), Inches(1.65), Inches(6.15), Inches(4.55))
    # Right: AI response
    img(sl, "chat_resp", Inches(6.75), Inches(1.65), Inches(6.15), Inches(4.55))

    # Caption strip below left
    rect(sl, Inches(0.4), Inches(6.28), Inches(6.15), Inches(0.42), RGBColor(0x09, 0x14, 0x2C))
    txt(sl, "Chat UI after login — sidebar shows session history",
        Inches(0.55), Inches(6.32), Inches(5.9), Inches(0.34),
        size=10, color=DIM, italic=True)

    # Caption strip below right
    rect(sl, Inches(6.75), Inches(6.28), Inches(6.15), Inches(0.42), RGBColor(0x09, 0x14, 0x2C))
    txt(sl, 'AI response to "Make me a classic Negroni cocktail" — real Gemini output',
        Inches(6.9), Inches(6.32), Inches(5.9), Inches(0.34),
        size=10, color=DIM, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 – Admin Panel
# ══════════════════════════════════════════════════════════════════════════════
def slide_admin(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl); particles(sl, 10, seed=44)
    header_band(sl, "Admin Panel & Security",
                "Role-based control — only admins can access system management")

    # Left column: two admin screenshots stacked
    img(sl, "admin",      Inches(0.4),  Inches(1.65), Inches(7.0), Inches(2.9))
    img(sl, "admin_chat", Inches(0.4),  Inches(4.65), Inches(7.0), Inches(2.1))

    # Right panel: feature breakdown
    rect(sl, Inches(7.6), Inches(1.65), Inches(5.3), Inches(5.1), PANEL2,
         RGBColor(0x22, 0x38, 0x60), Pt(0.6))
    vline(sl, Inches(7.6), Inches(1.65), Inches(5.1), GOLD)

    txt(sl, "Admin Capabilities", Inches(7.8), Inches(1.82), Inches(5.0), Inches(0.5),
        size=18, bold=True, color=GOLD_LT)
    hline(sl, Inches(7.8), Inches(2.38), Inches(4.8), GOLD, Inches(0.025))

    admin_items = [
        ("👤", "User Management",
         "Full roster of registered users — see ID, name, email, and role at a glance."),
        ("🔄", "Promote / Demote",
         "Toggle any user between User and Admin roles with a single button click."),
        ("🗑️", "Account Deletion",
         "Remove any account instantly. Built-in guard prevents admins from deleting themselves."),
        ("📜", "Recipe History",
         "Browse every AI-generated recipe across all users — user email, recipe name, timestamp."),
        ("💬", "Chat Log Viewer",
         "Read full conversation transcripts of any session. View button opens a modal overlay."),
    ]

    y = Inches(2.52)
    for icon, heading, detail in admin_items:
        dot_icon(sl, Inches(7.82), y + Inches(0.06), GOLD)
        txt(sl, icon + " " + heading, Inches(8.02), y, Inches(4.7), Inches(0.38),
            size=13, bold=True, color=WHITE)
        txt(sl, detail, Inches(8.02), y + Inches(0.37), Inches(4.7), Inches(0.38),
            size=10.5, color=MUTED, wrap=True)
        y += Inches(0.83)

    # Caption labels
    rect(sl, Inches(0.4), Inches(4.56), Inches(7.0), Inches(0.28), RGBColor(0x09, 0x14, 0x2C))
    txt(sl, "↑ Registered users table with role badges and action buttons",
        Inches(0.55), Inches(4.58), Inches(6.8), Inches(0.24),
        size=9, color=DIM, italic=True)

    rect(sl, Inches(0.4), Inches(6.75), Inches(7.0), Inches(0.28), RGBColor(0x09, 0x14, 0x2C))
    txt(sl, "↑ Global chat session history with user email, title, and view/delete controls",
        Inches(0.55), Inches(6.77), Inches(6.8), Inches(0.24),
        size=9, color=DIM, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 – Conclusion / Thank You
# ══════════════════════════════════════════════════════════════════════════════
def slide_conclusion(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl); particles(sl, 24, seed=3)

    # Large centred card
    rect(sl, Inches(0.8), Inches(0.45), Inches(11.75), Inches(6.6),
         RGBColor(0x0A, 0x16, 0x32), GOLD, Pt(0.8))
    hline(sl, Inches(0.8),  Inches(0.45),  Inches(11.75), GOLD, Inches(0.09))
    hline(sl, Inches(0.8),  Inches(7.05), Inches(11.75), GOLD, Inches(0.09))

    txt(sl, "🍹", Inches(5.85), Inches(0.9), Inches(1.6), Inches(1.05),
        size=52, align=PP_ALIGN.CENTER)

    txt(sl, "Thank You!", Inches(1.2), Inches(1.82), Inches(11.0), Inches(1.05),
        size=60, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    hline(sl, Inches(3.6), Inches(2.88), Inches(6.2), GOLD)

    txt(sl, "MixBot — Your AI-Powered Cocktail Mixologist",
        Inches(1.5), Inches(3.02), Inches(10.4), Inches(0.55),
        size=19, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

    # 2-column achievement checklist
    left_items = [
        "Full-stack web application (FastAPI + Vanilla JS)",
        "Gemini 2.0 Flash via OpenRouter integration",
        "JWT & Google OAuth authentication system",
    ]
    right_items = [
        "Persistent multi-session chat history (SQLite)",
        "Admin dashboard: users, recipes & chat logs",
        "Premium glassmorphism UI with particle animations",
    ]

    lx, rx = Inches(1.5), Inches(7.1)
    y = Inches(3.68)
    txt(sl, "Project Achievements", lx, y - Inches(0.38), Inches(5.2), Inches(0.38),
        size=12, bold=True, color=GOLD_LT)
    txt(sl, "Project Achievements", rx, y - Inches(0.38), Inches(5.2), Inches(0.38),
        size=12, bold=True, color=GOLD_LT)

    for item in left_items:
        dot_icon(sl, lx, y + Inches(0.08), GREEN)
        txt(sl, item, lx + Inches(0.22), y, Inches(5.3), Inches(0.38), size=12.5, color=WHITE)
        y += Inches(0.45)

    y = Inches(3.68)
    for item in right_items:
        dot_icon(sl, rx, y + Inches(0.08), GREEN)
        txt(sl, item, rx + Inches(0.22), y, Inches(5.3), Inches(0.38), size=12.5, color=WHITE)
        y += Inches(0.45)

    hline(sl, Inches(1.2), Inches(6.2), Inches(11.0), RGBColor(0x20, 0x30, 0x55))

    txt(sl, "Built with  Python · FastAPI · Google Gemini AI · HTML5 / CSS3 / JavaScript",
        Inches(1.2), Inches(6.28), Inches(11.0), Inches(0.4),
        size=11.5, color=DIM, align=PP_ALIGN.CENTER)

    txt(sl, "AI Essentials Capstone Project  |  April 2026",
        Inches(1.2), Inches(6.67), Inches(11.0), Inches(0.35),
        size=10.5, color=DIM, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
#  BUILD
# ══════════════════════════════════════════════════════════════════════════════
def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    steps = [
        ("Slide 1: Cover",              slide_cover),
        ("Slide 2: Project Overview",   slide_overview),
        ("Slide 3: Core Features",      slide_features),
        ("Slide 4: System Architecture",slide_architecture),
        ("Slide 5: Live Chat Interface",slide_chat),
        ("Slide 6: Admin Panel",        slide_admin),
        ("Slide 7: Conclusion",         slide_conclusion),
    ]

    for label, fn in steps:
        fn(prs)
        print(f"  [OK] {label}")

    prs.save(OUTPUT)
    print(f"\nSaved: {OUTPUT}")

if __name__ == "__main__":
    build()
